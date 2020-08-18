# Adding Images

import numpy as np
import pandas as pd
import plotly.graph_objects as go
import datetime
from pptx.util import Inches
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
from pptx import Presentation
from pptx.util import Inches


prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
# slide = prs.slides.add_slide(blank_slide_layout)

# Page 1
# -----------------------------------------------------------------------------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[6])

# shape = slide.shapes.add_shape(
#     MSO_SHAPE.RECTANGLE, 0, Inches(0.5), Inches(16), Inches(0.3))

left = top = width = height = Inches(1.0)

shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(0.5), Inches(10), Inches(0.5)
)




shape.shadow.inherit = False
fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 0, 0)

p = shape.text_frame.paragraphs[0]
run = p.add_run()
run.text = 'How to Add an Image?'

font = run.font
font.name = 'Calibri'
font.size = Pt(32)
font.bold = True
font.italic = None  # cause value to be inherited from theme
#font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

line = shape.line
line.color.rgb = RGBColor(255, 0, 0)

# logo1 = slide.shapes.add_picture(pylogo, Inches(
#     14.5), Inches(0.4), height=Inches(0.5), width=Inches(0.5))
# logo2 = slide.shapes.add_picture(pptlogo, Inches(
#     15.0), Inches(0.4), height=Inches(0.5), width=Inches(0.5))

left = Inches(0.5)
top = Inches(1.5)
width = Inches(9.5)
height = Inches(4)

text_box = slide.shapes.add_textbox(left, top, width, height)

tb = text_box.text_frame
tb.text = 'Do you know how the Orcs first came into being? They were elves once,\n\
taken by the dark powers, tortured and mutilated. A ruined and terrible form of life.\n\
Now... perfected. My fighting Uruk-Hai. Whom do you serve?'

prg = tb.add_paragraph()
prg.text = " "

prg = tb.add_paragraph()
prg.text = "They will find the Ring, and kill the one who carries it."

prs.save('font_001.pptx')
