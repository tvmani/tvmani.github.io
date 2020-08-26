from openpyxl import Workbook
from openpyxl.worksheet.table import Table, TableStyleInfo
from openpyxl import load_workbook
import collections as col
import string
import numpy as np
import pandas as pd
import plotly.graph_objects as go
import datetime
from pptx.util import Inches
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
from pptx import Presentation
from pptx.util import Inches
import os


# todo: filename should be passed in commandline for function
filename = "XIth_BIO-BOTANY_TM_Topics.xlsx"
# lessonTitle = "அலகு 5 - தாவர செயலியல் (Plant Physiology)"
lessonTitle = "11th - தாவரவியல் நீட் போட்டி தேர்வு பயிற்சி"

wb = load_workbook(filename)

sheet_ranges = wb.worksheets[0]

print(sheet_ranges.tables.items())

(table_name, table_cell_range) = sheet_ranges.tables.items()[0]

print(table_name, table_cell_range, sep=",")


def getTotalRecords(tableRange):
    # A1:C5
    (lower_cell, upper_cell) = tableRange.split(":")
    start_row = int(lower_cell[1:])
    end_row = int(upper_cell[1:])
    return end_row - start_row


def get_cell_identifier(tableRange):
    # A1:C5
    (lower_cell, upper_cell) = tableRange.split(":")
    start_row_id = lower_cell[0]
    end_row_id = upper_cell[0]
    return (start_row_id, end_row_id)


def get_cell_indexes(tableRange):
    # A1:C5
    (lower_cell, upper_cell) = tableRange.split(":")
    start_row_index = int(lower_cell[1:])
    end_row_index = int(upper_cell[1:]) - 1
    return (start_row_index, end_row_index)


def create_cover_slide(titleText, subtitleText, prs, isQuestion=True):
    lyt = prs.slide_layouts[0]  # choosing a slide layout

    slide = prs.slides.add_slide(lyt)  # adding a slide

    title = slide.shapes.title  # assigning a title

    fill = title.fill
    fill.solid()
    
    # fill.fore_color.rgb = RGBColor(255, 0, 0) # Blue
    fill.fore_color.rgb = RGBColor(0, 191, 255)  # Blue

    subtitle = slide.placeholders[1]  # placeholder for subtitle

    title.text = titleText
    # print(dir(title))

    subtitle.text = subtitleText  # subtitle


def getOptionAsString(optionInput):
    result = ""
    if optionInput is None:
        result = ""
    else:
        result = str(optionInput)
    return result


def create_slide(quizRecord, prs, isQuestion=True):

    blank_slide_layout = prs.slide_layouts[6]

    # Page 1
    # -----------------------------------------------------------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_slide_layout)

    left = top = width = height = Inches(1.0)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(0.1), Inches(10), Inches(0.6)
    )
    shape.shadow.inherit = False
    fill = shape.fill
    fill.solid()
    # fill.fore_color.rgb = RGBColor(255, 0, 0) # Blue
    fill.fore_color.rgb = RGBColor(0, 191, 255)  # Blue

    p = shape.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = lessonTitle + "\n" + quizRecord.topic_name

    font = run.font
    font.name = 'Calibri'
    font.size = Pt(20)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    # font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

    line = shape.line
    # line.color.rgb = RGBColor(255, 0, 0) # RED
    # line.color.rgb = RGBColor(0,191,255) # deep sky blue

    # logo1 = slide.shapes.add_picture(pylogo, Inches(
    #     14.5), Inches(0.4), height=Inches(0.5), width=Inches(0.5))
    # logo2 = slide.shapes.add_picture(pptlogo, Inches(
    #     15.0), Inches(0.4), height=Inches(0.5), width=Inches(0.5))

    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(8)
    height = Inches(4)

    text_box = slide.shapes.add_textbox(left, top, width, height)

    tb = text_box.text_frame
    tb.word_wrap = True

    currentQuiz = quizRecord
    # tb.text = "Q - " + str(currentQuiz.question_no) + \
    #     "  " + currentQuiz.question

    questionParagraph = tb.add_paragraph()
    questionParagraph.text = "Q - " + str(currentQuiz.question_no) + \
        "  " + currentQuiz.question

    font = questionParagraph.font
    font.name = 'Calibri'
    font.size = Pt(20)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    font.color.rgb = RGBColor(255, 0, 0)

    prg = tb.add_paragraph()
    prg.text = ""

    if isQuestion == True:
        prg = tb.add_paragraph()
        prg.text = "அ) " + getOptionAsString(currentQuiz.option_1)
        prg = tb.add_paragraph()
        prg.text = ""

        prg = tb.add_paragraph()
        prg.text = "ஆ) " + getOptionAsString(currentQuiz.option_2)
        prg = tb.add_paragraph()
        prg.text = ""

        prg = tb.add_paragraph()
        prg.text = "இ) " + getOptionAsString(currentQuiz.option_3)
        prg = tb.add_paragraph()
        prg.text = ""

        prg = tb.add_paragraph()
        prg.text = "ஈ) " + getOptionAsString(currentQuiz.option_4)

        if len(getOptionAsString(currentQuiz.option_5)) > 0:
            prg = tb.add_paragraph()
            prg.text = "உ) " + getOptionAsString(currentQuiz.option_5)
    else:
        prg = tb.add_paragraph()

        font = prg.font
        font.name = 'Calibri'
        font.size = Pt(20)
        font.bold = True
        font.italic = None  # cause value to be inherited from theme
        font.color.rgb = RGBColor(0, 100, 0)

        if quizRecord.correct_answer == 1:
            prg.text = "அ) " + getOptionAsString(currentQuiz.option_1)
        elif quizRecord.correct_answer == 2:
            prg.text = "ஆ) " + getOptionAsString(currentQuiz.option_2)
        elif quizRecord.correct_answer == 3:
            prg.text = "இ) " + getOptionAsString(currentQuiz.option_3)
        elif quizRecord.correct_answer == 4:
            prg.text = "ஈ) " + getOptionAsString(currentQuiz.option_4)
        elif quizRecord.correct_answer == 5:
            prg.text = "உ) " + getOptionAsString(currentQuiz.option_5)

    left = Inches(0.5)
    top = Inches(7.0)
    width = Inches(4)
    height = Inches(0.5)

    text_box = slide.shapes.add_textbox(left, top, width, height)

    tb = text_box.text_frame
    tb.word_wrap = True

    tb.text = "https://dailypractice.info/neet"


print(getTotalRecords(table_cell_range))
print(get_cell_identifier(table_cell_range))
print(get_cell_indexes(table_cell_range))
excel_A_to_Z = list(string.ascii_uppercase)

(start_row_index, end_row_index) = get_cell_indexes(table_cell_range)
(start_row_id, end_row_id) = get_cell_identifier(table_cell_range)

# question_no	question	exams_already_asked	year	option_1	option_2	option_3	option_4	correct_answer

QuizRecord = col.namedtuple('QuizRecord', 'question_no, question, exams_already_asked, year, option_1, \
                            option_2, option_3, option_4, option_5, correct_answer, page_no,topic_name')
QuizRecordsDb = []

for row in range(start_row_index, end_row_index, 1):
    print("processing row " + str(start_row_index))
    quizRecord = QuizRecord(sheet_ranges['A' + str(row + 1)].value,
                            sheet_ranges['B' + str(row + 1)].value,
                            sheet_ranges['C' + str(row + 1)].value,
                            sheet_ranges['D' + str(row + 1)].value,
                            sheet_ranges['E' + str(row + 1)].value,
                            sheet_ranges['F' + str(row + 1)].value,
                            sheet_ranges['G' + str(row + 1)].value,
                            sheet_ranges['H' + str(row + 1)].value,
                            sheet_ranges['I' + str(row + 1)].value,
                            sheet_ranges['J' + str(row + 1)].value,
                            sheet_ranges['K' + str(row + 1)].value,
                            sheet_ranges['L' + str(row + 1)].value,
                            )
    QuizRecordsDb.append(quizRecord)

prs = Presentation()
create_cover_slide(lessonTitle, "   ", prs)

previous_topic_name = ""
current_topic_name = ""

for quizRecord in QuizRecordsDb:
    current_topic_name = quizRecord.topic_name
    if current_topic_name != previous_topic_name :
        create_cover_slide(current_topic_name, "   ", prs)
        previous_topic_name = current_topic_name

    create_slide(quizRecord, prs, True)
    create_slide(quizRecord, prs, False)

(filenameWithoutExtn, fileExtn) = filename.split(".")

print(filenameWithoutExtn)
pptFilename = filenameWithoutExtn + ".pptx"
prs.save(pptFilename)

print("Created: " + pptFilename)
