from openpyxl import Workbook
from openpyxl.worksheet.table import Table, TableStyleInfo
from openpyxl import load_workbook
import collections as col
import string
import numpy as np
import pandas as pd
import plotly.graph_objects as go
import datetime
from pptx.util import Inches
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
from pptx import Presentation
from pptx.util import Inches
import os


# todo: filename should be passed in commandline for function
filename = "L-05.xlsx"

wb = load_workbook(filename)

sheet_ranges = wb.worksheets[0]

print(sheet_ranges.tables.items())

(table_name, table_cell_range) = sheet_ranges.tables.items()[0]

print(table_name, table_cell_range, sep=",")


def getTotalRecords(tableRange):
    # A1:C5
    (lower_cell, upper_cell) = tableRange.split(":")
    start_row = int(lower_cell[1:])
    end_row = int(upper_cell[1:])
    return end_row - start_row


def get_cell_identifier(tableRange):
    # A1:C5
    (lower_cell, upper_cell) = tableRange.split(":")
    start_row_id = lower_cell[0]
    end_row_id = upper_cell[0]
    return (start_row_id, end_row_id)


def get_cell_indexes(tableRange):
    # A1:C5
    (lower_cell, upper_cell) = tableRange.split(":")
    start_row_index = int(lower_cell[1:])
    end_row_index = int(upper_cell[1:])
    return (start_row_index, end_row_index)


def create_slide(quizRecord, prs, isQuestion=True):

    blank_slide_layout = prs.slide_layouts[6]

    # Page 1
    # -----------------------------------------------------------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_slide_layout)

    left = top = width = height = Inches(1.0)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(0.5), Inches(10), Inches(0.5)
    )
    shape.shadow.inherit = False
    fill = shape.fill
    fill.solid()
    #fill.fore_color.rgb = RGBColor(255, 0, 0) # RED
    fill.fore_color.rgb = RGBColor(0, 191, 255) # RED

    p = shape.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = '11th Std - Questions'

    font = run.font
    font.name = 'Calibri'
    font.size = Pt(24)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    # font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

    line = shape.line
    # line.color.rgb = RGBColor(255, 0, 0) # RED
    line.color.rgb = RGBColor(0,191,255) # deep sky blue    

    # logo1 = slide.shapes.add_picture(pylogo, Inches(
    #     14.5), Inches(0.4), height=Inches(0.5), width=Inches(0.5))
    # logo2 = slide.shapes.add_picture(pptlogo, Inches(
    #     15.0), Inches(0.4), height=Inches(0.5), width=Inches(0.5))

    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(4)

    text_box = slide.shapes.add_textbox(left, top, width, height)

    tb = text_box.text_frame
    tb.word_wrap = True

    currentQuiz = quizRecord
    tb.text = "Q - " + str(currentQuiz.question_no) + \
        "  " + currentQuiz.question

    # font = text_box.font
    # font.name = 'Calibri'
    # font.size = Pt(24)
    # font.bold = True

    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(4)

    text_box_options = slide.shapes.add_textbox(left, top, width, height)

    tbOptions = text_box_options.text_frame
    tbOptions.word_wrap = True

    prg = tbOptions.add_paragraph()
    prg.text = ""

    if isQuestion == True:
        prg = tbOptions.add_paragraph()
        prg.text = "அ) " + currentQuiz.option_1
        prg = tbOptions.add_paragraph()
        prg.text = ""

        prg = tbOptions.add_paragraph()
        prg.text = "ஆ) " + currentQuiz.option_2
        prg = tbOptions.add_paragraph()
        prg.text = ""

        prg = tbOptions.add_paragraph()
        prg.text = "இ) " + currentQuiz.option_3
        prg = tbOptions.add_paragraph()
        prg.text = ""

        prg = tbOptions.add_paragraph()
        prg.text = "ஈ) " + currentQuiz.option_4
    else:
        prg = tbOptions.add_paragraph()

        font = prg.font
        font.name = 'Calibri'
        font.size = Pt(24)
        font.bold = True
        font.italic = None  # cause value to be inherited from theme
        font.color.rgb = RGBColor(0, 100, 0)

        if quizRecord.correct_answer == 1:
            prg.text = "அ) " + currentQuiz.option_1
        elif quizRecord.correct_answer == 2:
            prg.text = "ஆ) " + currentQuiz.option_2
        elif quizRecord.correct_answer == 3:
            prg.text = "இ) " + currentQuiz.option_3
        elif quizRecord.correct_answer == 4:
            prg.text = "ஈ) " + currentQuiz.option_4

    left = Inches(0.5)
    top = Inches(7.0)
    width = Inches(4)
    height = Inches(0.5)

    text_box = slide.shapes.add_textbox(left, top, width, height)

    tb = text_box.text_frame
    tb.word_wrap = True

    tb.text = "https://dailypractice.info/neet"

print(getTotalRecords(table_cell_range))
print(get_cell_identifier(table_cell_range))
print(get_cell_indexes(table_cell_range))
excel_A_to_Z = list(string.ascii_uppercase)

(start_row_index, end_row_index) = get_cell_indexes(table_cell_range)
(start_row_id, end_row_id) = get_cell_identifier(table_cell_range)

# question_no	question	exams_already_asked	year	option_1	option_2	option_3	option_4	correct_answer

QuizRecord = col.namedtuple('QuizRecord', 'question_no, question, exams_already_asked, year, option_1, \
                            option_2, option_3, option_4, correct_answer')
QuizRecordsDb = []

for row in range(start_row_index, end_row_index, 1):
    print("processing row " + str(start_row_index))
    quizRecord = QuizRecord(sheet_ranges['A' + str(row + 1)].value,
                            sheet_ranges['B' + str(row + 1)].value,
                            sheet_ranges['C' + str(row + 1)].value,
                            sheet_ranges['D' + str(row + 1)].value,
                            sheet_ranges['E' + str(row + 1)].value,
                            sheet_ranges['F' + str(row + 1)].value,
                            sheet_ranges['G' + str(row + 1)].value,
                            sheet_ranges['H' + str(row + 1)].value,
                            sheet_ranges['I' + str(row + 1)].value,
                            )
    QuizRecordsDb.append(quizRecord)

prs = Presentation()

for quizRecord in QuizRecordsDb:
    # print(x)
    create_slide(quizRecord, prs, True)
    create_slide(quizRecord, prs, False)

(filenameWithoutExtn, fileExtn) = filename.split(".")

print(filenameWithoutExtn)
pptFilename = filenameWithoutExtn + ".pptx"
prs.save(pptFilename)

print("Created: " + pptFilename)
