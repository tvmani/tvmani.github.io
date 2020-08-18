# Adding Images

from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

path = r'images\620px-Duloxetine-ball-and-stick-2.png'

left = Inches(1)
top = Inches(0.5)

img = slide.shapes.add_picture(path, left, top)

prs.save("slide2.pptx")  # saving file
